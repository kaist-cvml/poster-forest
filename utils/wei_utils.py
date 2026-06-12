import re
import io
import contextlib
import traceback
import threading
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE, MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from camel.types import ModelPlatformType, ModelType
from camel.configs import ChatGPTConfig, QwenConfig, VLLMConfig, OpenRouterConfig, GeminiConfig
import math
from urllib.parse import quote_from_bytes, quote
from PIL import Image
import os
import copy
from utils.src.utils import ppt_to_images
from playwright.sync_api import sync_playwright
from pathlib import Path
from playwright.async_api import async_playwright
import asyncio
from utils.pptx_utils import *
from utils.critic_utils import *
from camel.models import ModelFactory

# ---------------------------------------------------------------------------
# Agent step with timeout + retry
# ---------------------------------------------------------------------------
_AGENT_STEP_TIMEOUT_VISION = 180   # seconds — vision calls (image_list)
_AGENT_STEP_TIMEOUT_TEXT   = 120   # seconds — text-only calls
_AGENT_STEP_MAX_RETRIES    = 2

def agent_step_with_timeout(agent, msg, timeout=None, max_retries=_AGENT_STEP_MAX_RETRIES):
    """Call agent.step(msg) with a wall-clock timeout and automatic retry.

    Uses daemon=True threads so a stuck network call never blocks the main
    thread or prevents the process from exiting.  thread.join(timeout) returns
    immediately when the deadline expires regardless of what the thread is doing.

    timeout defaults to _AGENT_STEP_TIMEOUT_TEXT (120 s) for text messages
    and _AGENT_STEP_TIMEOUT_VISION (180 s) when msg carries images.
    """
    if timeout is None:
        has_images = hasattr(msg, 'image_list') and bool(msg.image_list)
        timeout = _AGENT_STEP_TIMEOUT_VISION if has_images else _AGENT_STEP_TIMEOUT_TEXT

    last_exc = None
    for attempt in range(max_retries + 1):
        result_box = [None]
        exc_box    = [None]

        def _run(result_box=result_box, exc_box=exc_box):
            try:
                result_box[0] = agent.step(msg)
            except Exception as e:
                exc_box[0] = e

        t = threading.Thread(target=_run, daemon=True)
        t.start()
        t.join(timeout=timeout)

        if t.is_alive():
            # Thread is still blocked on the network — timeout exceeded.
            last_exc = TimeoutError(
                f"agent.step() timed out after {timeout}s "
                f"(attempt {attempt + 1}/{max_retries + 1})"
            )
            tag = 'retrying...' if attempt < max_retries else 'giving up.'
            print(f"   ⏰ API call timed out after {timeout}s "
                  f"(attempt {attempt + 1}/{max_retries + 1}) — {tag}")
            continue

        if exc_box[0] is not None:
            last_exc = exc_box[0]
            tag = 'retrying...' if attempt < max_retries else 'giving up.'
            print(f"   ❌ API call error: {exc_box[0]} "
                  f"(attempt {attempt + 1}/{max_retries + 1}) — {tag}")
            continue

        return result_box[0]

    raise last_exc


def make_model(config):
    """Create a CAMEL ModelFactory instance from an agent_config dict.

    Handles both vLLM (url key present) and cloud models uniformly,
    so callers don't need to repeat the if/else startswith('vllm_qwen') pattern.
    """
    kwargs = {
        'model_platform': config['model_platform'],
        'model_type': config['model_type'],
        'model_config_dict': config['model_config'],
    }
    if 'url' in config:
        kwargs['url'] = config['url']
    return ModelFactory.create(**kwargs)


# model_name_t / model_name_v  →  model                              endpoint
# vllm_qwen3                      Qwen/Qwen3-30B-A3B-Instruct-2507   localhost:8005
# vllm_qwen3_vl                   Qwen/Qwen3-VL-8B-Instruct          localhost:8010
# vllm_qwen2_5                    Qwen/Qwen2.5-7B-Instruct           localhost:8005
# vllm_qwen2_5_vl                 Qwen/Qwen2.5-VL-7B-Instruct        localhost:8010
# 4o                              gpt-4o                             OpenAI API
# 4o-mini                         gpt-4o-mini                        OpenAI API
def get_agent_config(model_type):
    agent_config = {}
    if model_type == 'qwen':
        agent_config = {
            "model_type": ModelType.DEEPINFRA_QWEN_2_5_72B,
            "model_config": QwenConfig().as_dict(),
            "model_platform": ModelPlatformType.DEEPINFRA,
        }
    elif model_type == 'gemini':
        agent_config = {
            "model_type": ModelType.DEEPINFRA_GEMINI_2_FLASH,
            "model_config": GeminiConfig().as_dict(),
            "model_platform": ModelPlatformType.DEEPINFRA,
            'max_images': 99
        }
    elif model_type == 'phi4':
        agent_config = {
            "model_type": ModelType.DEEPINFRA_PHI_4_MULTIMODAL,
            "model_config": QwenConfig().as_dict(),
            "model_platform": ModelPlatformType.DEEPINFRA,
        }
    elif model_type == 'llama-4-scout-17b-16e-instruct':
        agent_config = {
            'model_type': ModelType.ALIYUN_LLAMA4_SCOUT_17B_16E,
            'model_config': QwenConfig().as_dict(),
            'model_platform': ModelPlatformType.QWEN,
            'max_images': 99
        }
    elif model_type == 'qwen-2.5-vl-72b':
        agent_config = {
            'model_type': ModelType.QWEN_2_5_VL_72B,
            'model_config': QwenConfig().as_dict(),
            'model_platform': ModelPlatformType.QWEN,
            'max_images': 99
        }
    elif model_type == 'gemma':
        agent_config = {
            "model_type": "google/gemma-3-4b-it",
            "model_platform": ModelPlatformType.VLLM,
            "model_config": VLLMConfig().as_dict(),
            "url": 'http://localhost:5555/v1',
            'max_images': 99
        }
    elif model_type == 'llava':
        agent_config = {
            "model_type": "llava-hf/llava-onevision-qwen2-7b-ov-hf",
            "model_platform": ModelPlatformType.VLLM,
            "model_config": VLLMConfig().as_dict(),
            "url": 'http://localhost:8000/v1',
            'max_images': 99
        }
    elif model_type == 'molmo-o':
        agent_config = {
            "model_type": "allenai/Molmo-7B-O-0924",
            "model_platform": ModelPlatformType.VLLM,
            "model_config": VLLMConfig().as_dict(),
            "url": 'http://localhost:8000/v1',
            'max_images': 99
        }
    elif model_type == 'qwen-2-vl-7b':
        agent_config = {
            "model_type": "Qwen/Qwen2-VL-7B-Instruct",
            "model_platform": ModelPlatformType.VLLM,
            "model_config": VLLMConfig().as_dict(),
            "url": 'http://localhost:8000/v1',
            'max_images': 99
        }
    elif model_type == 'vllm_phi4':
        agent_config = {
            "model_type": "microsoft/Phi-4-multimodal-instruct",
            "model_platform": ModelPlatformType.VLLM,
            "model_config": VLLMConfig().as_dict(),
            "url": 'http://localhost:8000/v1',
            'max_images': 99
        }
    elif model_type == 'o3-mini':
        agent_config = {
            "model_type": ModelType.O3_MINI,
            "model_config": ChatGPTConfig().as_dict(),
            "model_platform": ModelPlatformType.OPENAI,
        }
    elif model_type == 'gpt-4.1':
        agent_config = {
            "model_type": ModelType.GPT_4_1,
            "model_config": ChatGPTConfig().as_dict(),
            "model_platform": ModelPlatformType.OPENAI,
        }
    elif model_type == 'gpt-4.1-mini':
        agent_config = {
            "model_type": ModelType.GPT_4_1_MINI,
            "model_config": ChatGPTConfig().as_dict(),
            "model_platform": ModelPlatformType.OPENAI,
        }
    elif model_type == '4o':
        agent_config = {
            "model_type": ModelType.GPT_4O,
            "model_config": ChatGPTConfig().as_dict(),
            "model_platform": ModelPlatformType.OPENAI,
            "token_limit": 32768,
        }
    elif model_type == '4o-mini':
        agent_config = {
            "model_type": ModelType.GPT_4O_MINI,
            "model_config": ChatGPTConfig().as_dict(),
            "model_platform": ModelPlatformType.OPENAI,
        }
    elif model_type == 'o1':
        agent_config = {
            "model_type": ModelType.O1,
            "model_config": ChatGPTConfig().as_dict(),
            "model_platform": ModelPlatformType.OPENAI,
        }
    elif model_type == 'o3':
        agent_config = {
            "model_type": ModelType.O3,
            "model_config": ChatGPTConfig().as_dict(),
            "model_platform": ModelPlatformType.OPENAI,
        }
    elif model_type == 'vllm_qwen3_vl':
        agent_config = {
            "model_type": "Qwen/Qwen3-VL-8B-Instruct",
            "model_platform": ModelPlatformType.VLLM,
            "model_config": VLLMConfig().as_dict(),
            "url": 'http://localhost:8010/v1',
            'max_images': 99
        }
    elif model_type == 'vllm_qwen3':
        agent_config = {
            "model_type": "Qwen/Qwen3-30B-A3B-Instruct-2507",
            "model_platform": ModelPlatformType.VLLM,
            "model_config": VLLMConfig().as_dict(),
            "url": 'http://localhost:8005/v1',
        }
    elif model_type == 'vllm_qwen2_5':
        agent_config = {
            "model_type": "Qwen/Qwen2.5-7B-Instruct",
            "model_platform": ModelPlatformType.VLLM,
            "model_config": VLLMConfig().as_dict(),
            "url": 'http://localhost:8005/v1',
        }
    elif model_type == 'vllm_qwen2_5_vl':
        agent_config = {
            "model_type": "Qwen/Qwen2.5-VL-7B-Instruct",
            "model_platform": ModelPlatformType.VLLM,
            "model_config": VLLMConfig().as_dict(),
            "url": 'http://localhost:8010/v1',
            'max_images': 99
        }
    elif model_type == 'openrouter_qwen_vl_72b':
        agent_config = {
            'model_type': ModelType.OPENROUTER_QWEN_2_5_VL_72B,
            'model_platform': ModelPlatformType.OPENROUTER,
            'model_config': OpenRouterConfig().as_dict(),
        }
    elif model_type == 'openrouter_qwen_vl_7b':
        agent_config = {
            'model_type': ModelType.OPENROUTER_QWEN_2_5_VL_7B,
            'model_platform': ModelPlatformType.OPENROUTER,
            'model_config': OpenRouterConfig().as_dict(),
        }
    elif model_type == 'openrouter_qwen_7b':
        agent_config = {
            'model_type': ModelType.OPENROUTER_QWEN_2_5_7B,
            'model_platform': ModelPlatformType.OPENROUTER,
            'model_config': OpenRouterConfig().as_dict(),
        }
    else:
        agent_config = {
            'model_type': model_type,
            'model_platform': ModelPlatformType.OPENAI_COMPATIBLE_MODEL,
            'model_config': None
        }
    
    return agent_config


def match_response(response):
    response_text = response.msgs[0].content
    pattern = r'```python(.*?)```'
    match = re.search(pattern, response_text, flags=re.DOTALL)
    if not match:
        pattern = r'```(.*?)```'
        match = re.search(pattern, response_text, flags=re.DOTALL)
    return match.group(1).strip() if match else response_text

def run_code_with_utils(code, utils_functions):
    return run_code(utils_functions + '\n' + code)

def run_code(code):
    """Execute code, capturing stdout. Returns (output, traceback_or_None).

    Sets __name__ = '__main__' so if-main blocks execute.
    """
    stdout_capture = io.StringIO()
    exec_globals = {"__name__": "__main__"}
    with contextlib.redirect_stdout(stdout_capture):
        try:
            exec(code, exec_globals)
            error = None
        except Exception:
            error = traceback.format_exc()
    return stdout_capture.getvalue(), error


def run_code_from_agent(agent, msg, num_retries=1):
    agent.reset()
    log = []
    for attempt in range(num_retries + 1):  # +1 to include the initial attempt
        response = agent.step(msg)
        code = match_response(response)
        output, error = run_code(code)
        log.append((code, output, error))
        
        if error is None:
            return log
        
        if attempt < num_retries:
            print(f"Retrying... Attempt {attempt + 1} of {num_retries}")
            msg = error
    
    return log

def run_modular(all_code, file_name, with_border=True, with_label=True):
    concatenated_code = utils_functions
    concatenated_code += "\n".join(all_code.values())
    if with_border and with_label:
        concatenated_code += add_border_label_function
        concatenated_code += create_id_map_function
        concatenated_code += save_helper_info_border_label.format(file_name, file_name, file_name)
    elif with_border:
        concatenated_code += add_border_function
        concatenated_code += save_helper_info_border.format(file_name, file_name)
    else:
        concatenated_code += f'\nposter.save("{file_name}")'
    output, error = run_code(concatenated_code)
    return concatenated_code, output, error

def edit_modular(
        agent,
        edit_section_name, 
        feedback,
        all_code, 
        file_name, 
        outline,
        content,
        images,
        actor_prompt,
        num_retries=1,
        prompt_type='initial'
    ):
    agent.reset()
    log = []
    if prompt_type == 'initial':
        msg = actor_prompt.format(
            outline['meta'],
            {edit_section_name: outline[edit_section_name]}, 
            content, 
            images,
            documentation
        )
    elif prompt_type == 'edit':
        assert (edit_section_name == list(feedback.keys())[0])
        msg = actor_prompt.format(
            edit_section_name,
            all_code[edit_section_name],
            feedback,
            {edit_section_name: outline[edit_section_name]}, 
            content, 
            images,
            documentation
        )
    elif prompt_type == 'new':
        assert (list(feedback.keys())[0] == 'all_good')
        msg = actor_prompt.format(
            {edit_section_name: outline[edit_section_name]}, 
            content, 
            images,
            documentation
        )

    for attempt in range(num_retries + 1):
        response = agent.step(msg)
        new_code = match_response(response)
        all_code_changed = all_code.copy()
        all_code_changed[edit_section_name] = new_code
        concatenated_code, output, error = run_modular(all_code_changed, file_name, False, False)
        log.append({
            "code": new_code,
            "output": output,
            "error": error,
            "concatenated_code": concatenated_code
        })
        if error is None:
            return log
        
        if attempt < num_retries:
            print(f"Retrying... Attempt {attempt + 1} of {num_retries}")
            msg = error
            msg += '\nFix your code and try again. The poster is a single-page pptx.'
            if prompt_type != 'initial':
                msg += '\nAssume that you have had a Presentation object named "poster" and a slide named "slide".'

    return log

def add_border_to_all_elements(prs, border_color=RGBColor(255, 0, 0), border_width=Pt(2)):
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                shape.line.fill.solid()
                shape.line.fill.fore_color.rgb = border_color
                shape.line.width = border_width
            except Exception as e:
                print(f"Could not add border to shape {shape.shape_type}: {e}")


# 1 point = 12700 EMUs (helper function)
def pt_to_emu(points: float) -> int:
    return int(points * 12700)

def add_border_and_labels(
    prs,
    border_color=RGBColor(255, 0, 0),
    border_width=Pt(2),
    label_outline_color=RGBColor(0, 0, 255),
    label_text_color=RGBColor(0, 0, 255),
    label_diameter_pt=40
):
    """Add red border + numbered blue circle label to each shape for debugging."""
    label_diameter_emu = pt_to_emu(label_diameter_pt)
    label_counter = 0
    labeled_elements = {}

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.name.startswith("Label_"):
                continue
            try:
                shape.line.fill.solid()
                shape.line.fill.fore_color.rgb = border_color
                shape.line.width = border_width

                label_left = shape.left + (shape.width // 2) - (label_diameter_emu // 2)
                label_top  = shape.top  + (shape.height // 2) - (label_diameter_emu // 2)
                label_shape = slide.shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.OVAL, label_left, label_top,
                    label_diameter_emu, label_diameter_emu
                )
                label_shape.name = f"Label_{label_counter}"
                label_shape.fill.background()
                label_shape.line.fill.solid()
                label_shape.line.fill.fore_color.rgb = label_outline_color
                label_shape.line.width = Pt(3)

                tf = label_shape.text_frame
                tf.text = str(label_counter)
                paragraph = tf.paragraphs[0]
                paragraph.alignment = PP_ALIGN.CENTER
                run = paragraph.runs[0]
                font = run.font
                font.size = Pt(40)
                font.bold = True
                font.name = "Arial"
                font._element.get_or_change_to_solidFill()
                font.fill.fore_color.rgb = label_text_color

                labeled_elements[label_counter] = {
                    'left': f'{shape.left} EMU',
                    'top': f'{shape.top} EMU',
                    'width': f'{shape.width} EMU',
                    'height': f'{shape.height} EMU',
                    'font_size': f'{shape.text_frame.font.size} PT' if hasattr(shape, 'text_frame') else None,
                }
                label_counter += 1
            except Exception as e:
                print(f"Could not add border/label to shape (type={shape.shape_type}): {e}")

    return labeled_elements


def fill_content(agent, prompt, num_retries, existing_code=''):
    if existing_code == '':
        existing_code = utils_functions
    agent.reset()
    log = []
    cumulative_input_token, cumulative_output_token = 0, 0
    for attempt in range(num_retries + 1):
        response = agent.step(prompt)
        input_token, output_token = account_token(response)
        cumulative_input_token += input_token
        cumulative_output_token += output_token
        new_code = match_response(response)
        all_code = existing_code + '\n' + new_code

        output, error = run_code(all_code)
        log.append({
            "code": new_code,
            "output": output,
            "error": error,
            "concatenated_code": all_code,
            'cumulative_tokens': (cumulative_input_token, cumulative_output_token)
        })

        if error is None:
            return log
        
        if attempt < num_retries:
            print(f"Retrying... Attempt {attempt + 1} of {num_retries}")
            prompt = error
    return log

def apply_theme(agent, prompt, num_retries, existing_code=''):
    return fill_content(agent, prompt, num_retries, existing_code)

def edit_code(agent, prompt, num_retries, existing_code=''):
    return fill_content(agent, prompt, num_retries, existing_code)

def stylize(agent, prompt, num_retries, existing_code=''):
    return fill_content(agent, prompt, num_retries, existing_code)

def gen_layout(agent, prompt, num_retries, name_to_hierarchy, visual_identifier='', existing_code=''):
    if existing_code == '':
        existing_code = utils_functions
    agent.reset()
    log = []
    cumulative_input_token, cumulative_output_token = 0, 0
    for attempt in range(num_retries + 1):
        response = agent.step(prompt)
        input_token, output_token = account_token(response)
        cumulative_input_token += input_token
        cumulative_output_token += output_token
        new_code = match_response(response)
        all_code = existing_code + '\n' + new_code

        # Save visualizations
        all_code += f'''
name_to_hierarchy = {name_to_hierarchy}
identifier = "{visual_identifier}"
get_visual_cues(name_to_hierarchy, identifier)
'''

        output, error = run_code(all_code)
        log.append({
            "code": new_code,
            "output": output,
            "error": error,
            "concatenated_code": all_code,
            'num_tokens': (input_token, output_token),
            'cumulative_tokens': (cumulative_input_token, cumulative_output_token)
        })

        if error is None:
            return log
        
        if attempt < num_retries:
            print(f"Retrying... Attempt {attempt + 1} of {num_retries}")
            prompt = error
    return log

def gen_layout_parallel(agent, prompt, num_retries, existing_code='', slide_width=0, slide_height=0, tmp_name='tmp'):
    if existing_code == '':
        existing_code = utils_functions
        
    existing_code += f'''
poster = create_poster(width_inch={slide_width}, height_inch={slide_height})
slide = add_blank_slide(poster)
save_presentation(poster, file_name="poster_{tmp_name}.pptx")
'''
    agent.reset()
    log = []
    cumulative_input_token, cumulative_output_token = 0, 0
    for attempt in range(num_retries + 1):
        response = agent.step(prompt)
        input_token, output_token = account_token(response)
        cumulative_input_token += input_token
        cumulative_output_token += output_token
        new_code = match_response(response)
        all_code = existing_code + '\n' + new_code

        output, error = run_code(all_code)
        log.append({
            "code": new_code,
            "output": output,
            "error": error,
            "concatenated_code": all_code,
            'num_tokens': (input_token, output_token),
            'cumulative_tokens': (cumulative_input_token, cumulative_output_token)
        })
        if output is None or output == '':
            prompt = 'No object name printed.'
            continue

        if error is None:
            return log
        
        if attempt < num_retries:
            # print(f"Retrying... Attempt {attempt + 1} of {num_retries}", flush=True)
            prompt = error
    return log

def compute_bullet_length(textbox_content):
    total = 0
    for bullet in textbox_content:
        for run in bullet['runs']:
            total += len(run['text'])
    return total

def check_bounding_boxes(bboxes, overall_width, overall_height):
    """Check for overlaps and overflows in bounding boxes.

    Returns (box1, box2) on overlap, (box,) on overflow, () if clean.
    """
    box_list = []
    for name, coords in bboxes.items():
        box_list.append((name, coords["left"], coords["top"], coords["width"], coords["height"]))

    def boxes_overlap(box_a, box_b):
        _, left_a, top_a, width_a, height_a = box_a
        _, left_b, top_b, width_b, height_b = box_b
        no_overlap = (
            left_a + width_a <= left_b or left_b + width_b <= left_a or
            top_a + height_a <= top_b or top_b + height_b <= top_a
        )
        return not no_overlap

    n = len(box_list)
    for i in range(n):
        for j in range(i + 1, n):
            if boxes_overlap(box_list[i], box_list[j]):
                return (box_list[i][0], box_list[j][0])

    for name, left, top, width, height in box_list:
        if left < 0 or top < 0 or left + width > overall_width or top + height > overall_height:
            return (name,)

    return ()


def is_poster_filled(
    bounding_boxes: dict,
    overall_width: float,
    overall_height: float,
    max_lr_margin: float,
    max_tb_margin: float
) -> bool:
    """Return True if bounding boxes collectively fill the poster within margin constraints."""
    if not bounding_boxes:
        return False

    min_left   = min(b["left"] for b in bounding_boxes.values())
    max_right  = max(b["left"] + b["width"] for b in bounding_boxes.values())
    min_top    = min(b["top"] for b in bounding_boxes.values())
    max_bottom = max(b["top"] + b["height"] for b in bounding_boxes.values())

    if (min_left > max_lr_margin or overall_width - max_right > max_lr_margin or
            min_top > max_tb_margin or overall_height - max_bottom > max_tb_margin):
        return False
    return True

def check_and_fix_subsections(section, subsections):
    """Validate subsections fit within section and greedy-expand them if they don't fill it.

    Returns a tuple of offending names on boundary/overlap violation,
    an expanded dict if area is underfilled, or () if all is well.
    """
    def right(rect):
        return rect["left"] + rect["width"]

    def bottom(rect):
        return rect["top"] + rect["height"]

    def is_overlapping(r1, r2):
        return not (
            right(r1) <= r2["left"]
            or r1["left"] >= right(r2)
            or bottom(r1) <= r2["top"]
            or r1["top"] >= bottom(r2)
        )

    names_violating = set()
    sec_left, sec_top = section["left"], section["top"]
    sec_right = section["left"] + section["width"]
    sec_bottom = section["top"] + section["height"]

    for name, sub in subsections.items():
        if (sub["left"] < sec_left or sub["top"] < sec_top or
                right(sub) > sec_right or bottom(sub) > sec_bottom):
            names_violating.add(name)

    sub_keys = list(subsections.keys())
    for i in range(len(sub_keys)):
        for j in range(i + 1, len(sub_keys)):
            n1, n2 = sub_keys[i], sub_keys[j]
            if is_overlapping(subsections[n1], subsections[n2]):
                names_violating.add(n1)
                names_violating.add(n2)

    if names_violating:
        return tuple(sorted(names_violating))

    area_section = section["width"] * section["height"]
    area_subs = sum(sub["width"] * sub["height"] for sub in subsections.values())

    if area_subs < area_section:
        expanded_subs = {
            name: dict(sub) for name, sub in subsections.items()
        }

        def touching_left(sname, sbox):
            if abs(sbox["left"] - sec_left) < 1e-9:
                return True
            for oname, obox in expanded_subs.items():
                if oname == sname:
                    continue
                if abs(right(obox) - sbox["left"]) < 1e-9:
                    return True
            return False

        def touching_right(sname, sbox):
            r = right(sbox)
            if abs(r - sec_right) < 1e-9:
                return True
            for oname, obox in expanded_subs.items():
                if oname != sname and abs(obox["left"] - r) < 1e-9:
                    return True
            return False

        def touching_top(sname, sbox):
            if abs(sbox["top"] - sec_top) < 1e-9:
                return True
            for oname, obox in expanded_subs.items():
                if oname != sname and abs(bottom(obox) - sbox["top"]) < 1e-9:
                    return True
            return False

        def touching_bottom(sname, sbox):
            b = bottom(sbox)
            if abs(b - sec_bottom) < 1e-9:
                return True
            for oname, obox in expanded_subs.items():
                if oname != sname and abs(obox["top"] - b) < 1e-9:
                    return True
            return False

        for name in expanded_subs:
            sub = expanded_subs[name]

            if not touching_left(name, sub):
                left_bound = sec_left
                for oname, obox in expanded_subs.items():
                    if oname == name:
                        continue
                    r_ = obox["left"] + obox["width"]
                    if r_ <= sub["left"] and r_ > left_bound:
                        left_bound = r_
                delta = sub["left"] - left_bound
                if delta > 1e-9:
                    sub["width"] += delta
                    sub["left"] = left_bound

            if not touching_right(name, sub):
                right_bound = sec_right
                sub_right = sub["left"] + sub["width"]
                for oname, obox in expanded_subs.items():
                    if oname == name:
                        continue
                    left_ = obox["left"]
                    if left_ >= sub_right and left_ < right_bound:
                        right_bound = left_
                delta = right_bound - (sub["left"] + sub["width"])
                if delta > 1e-9:
                    sub["width"] += delta

            if not touching_top(name, sub):
                top_bound = sec_top
                for oname, obox in expanded_subs.items():
                    if oname == name:
                        continue
                    b_ = obox["top"] + obox["height"]
                    if b_ <= sub["top"] and b_ > top_bound:
                        top_bound = b_
                delta = sub["top"] - top_bound
                if delta > 1e-9:
                    sub["height"] += delta
                    sub["top"] = top_bound

            if not touching_bottom(name, sub):
                bottom_bound = sec_bottom
                sub_bottom = sub["top"] + sub["height"]
                for oname, obox in expanded_subs.items():
                    if oname == name:
                        continue
                    other_top = obox["top"]
                    if other_top >= sub_bottom and other_top < bottom_bound:
                        bottom_bound = other_top
                delta = bottom_bound - (sub["top"] + sub["height"])
                if delta > 1e-9:
                    sub["height"] += delta

        return expanded_subs

    return ()

async def rendered_dims(html: Path) -> tuple[int, int]:
    async with async_playwright() as p:
        browser = await p.chromium.launch()
        page    = await browser.new_page()        # no fixed viewport yet
        resolved = html.resolve()
        # quote_from_bytes expects bytes, so we encode the path as UTF‐8:
        url = "file://" + quote_from_bytes(str(resolved).encode("utf-8"), safe="/:")
        await page.goto(url, wait_until="networkidle")

        # 1) bounding-box of <body>
        body_box = await page.eval_on_selector(
            "body",
            "el => el.getBoundingClientRect()")
        w = int(body_box["width"])
        h = int(body_box["height"])

        await browser.close()
        return w, h

    
def html_to_png(html_abs_path, poster_width_default, poster_height_default, output_path):
    html_file = html_abs_path

    try:
        w, h = asyncio.run(rendered_dims(html_file))
    except:
        w = poster_width_default
        h = poster_height_default

    with sync_playwright() as p:
        path_posix = Path(html_file).resolve().as_posix()

        file_url = "file://" + quote(path_posix, safe="/:")
        browser = p.chromium.launch()
        page    = browser.new_page(viewport={"width": w, "height": h})
        page.goto(file_url, wait_until='networkidle')
        page.screenshot(path=output_path, full_page=True)
        browser.close()

def account_token(response):
    input_token = response.info['usage']['prompt_tokens']
    output_token = response.info['usage']['completion_tokens']

    return input_token, output_token

def style_bullet_content(bullet_content_item, color, fill_color):
    for i in range(len(bullet_content_item)):
        bullet_content_item[i]['runs'][0]['color'] = color
        bullet_content_item[i]['runs'][0]['fill_color'] = fill_color

def scale_to_target_area(width, height, target_width=900, target_height=1200):
    """Scale (width, height) uniformly so the area matches target_width * target_height."""
    scale_factor = math.sqrt(target_width * target_height / (width * height))
    return int(round(width * scale_factor)), int(round(height * scale_factor))

def char_capacity(
    bbox,
    font_size_px=40 * (96 / 72),
    *,
    avg_width_ratio: float = 0.54,
    line_height_ratio: float = 1,
    padding_px: int = 0,
) -> int:
    """Estimate character capacity of a text box given its pixel dimensions."""
    CHAR_CONST = 10
    _, _, height_px, width_px = bbox

    usable_w = max(0, width_px - 2 * padding_px)
    usable_h = max(0, height_px - 2 * padding_px)
    if usable_w == 0 or usable_h == 0:
        return 0

    chars_per_line = max(1, math.floor(usable_w / (font_size_px * avg_width_ratio)))
    lines = max(1, math.floor(usable_h / (font_size_px * line_height_ratio)))
    return chars_per_line * lines * CHAR_CONST

def estimate_characters(width_in_inches, height_in_inches, font_size_points, line_spacing_points=None):
    """Estimate number of characters that fit in an inch-dimensioned bounding box."""
    if line_spacing_points is None:
        line_spacing_points = 1.5 * font_size_points
    width_in_points = width_in_inches * 72
    height_in_points = height_in_inches * 72
    chars_per_line = int(width_in_points // (0.5 * font_size_points))
    lines_count = int(height_in_points // line_spacing_points)
    return chars_per_line * lines_count

def equivalent_length_with_forced_breaks(text, width_in_inches, font_size_points):
    """Return equivalent character length accounting for forced newlines consuming full-line slots."""
    max_chars_per_line = int(width_in_inches * 72 // (0.5 * font_size_points))
    total_equiv_length = 0
    for line in text.split('\n'):
        if not line:
            total_equiv_length += max_chars_per_line
        else:
            total_equiv_length += math.ceil(len(line) / max_chars_per_line) * max_chars_per_line
    return total_equiv_length

def actual_rendered_length(
    text,
    width_in_inches,
    height_in_inches,
    font_size_points,
    line_spacing_points=None
):
    """Estimate how many characters from text actually render in the bounding box."""
    if line_spacing_points is None:
        line_spacing_points = 1.5 * font_size_points

    max_chars_per_line = int(width_in_inches * 72 // (0.5 * font_size_points))
    max_lines = int(height_in_inches * 72 // line_spacing_points)

    used_lines = 0
    displayed_chars = 0

    for line in text.split('\n'):
        if not line:
            used_lines += 1
            if used_lines >= max_lines:
                break
            continue

        sub_lines = math.ceil(len(line) / max_chars_per_line)
        if used_lines + sub_lines <= max_lines:
            displayed_chars += len(line)
            used_lines += sub_lines
        else:
            lines_left = max_lines - used_lines
            if lines_left <= 0:
                break
            displayed_chars += min(lines_left * max_chars_per_line, len(line))
            break

    return displayed_chars


def remove_hierarchy_and_id(data):
    """Recursively strip 'hierarchy', 'id', and 'location' keys from a nested dict."""
    if isinstance(data, dict):
        return {
            k: remove_hierarchy_and_id(v)
            for k, v in data.items()
            if k not in ("hierarchy", "id", "location")
        }
    if isinstance(data, list):
        return [remove_hierarchy_and_id(item) for item in data]
    return data
    
def outline_estimate_num_chars(outline):
    for k, v in outline.items():
        if k == 'meta':
            continue
        if 'title' in k.lower() or 'author' in k.lower() or 'reference' in k.lower():
            continue
        if not 'subsections' in v:
            num_chars = estimate_characters(
                v['location']['width'], 
                v['location']['height'], 
                60, line_spacing_points=None
            )
            v['num_chars'] = num_chars
        else:
            for k_sub, v_sub in v['subsections'].items():
                if 'title' in k_sub.lower():
                    continue
                if 'path' in v_sub:
                    continue
                num_chars = estimate_characters(
                    v_sub['location']['width'], 
                    v_sub['location']['height'], 
                    60, line_spacing_points=None
                )
                v_sub['num_chars'] = num_chars

def generate_length_suggestions(result_json, original_section_outline, raw_section_outline):
    NOT_CHANGE = 'Do not change text.'
    original_section_outline = json.loads(original_section_outline)
    suggestion_flag = False
    new_section_outline = copy.deepcopy(result_json)
    def check_length(text, target, width, height):
        text_length = equivalent_length_with_forced_breaks(
            text,
            width,
            font_size_points=60,
        )
        if text_length - target > 100:
            return f'Text too long, shrink by {text_length - target} characters.'
        elif target - text_length > 100:
            return f'Text too short, expand by {target - text_length} characters.'
        else:
            return NOT_CHANGE

    if 'num_chars' in original_section_outline:
        new_section_outline['suggestions'] = check_length(
            result_json['description'], 
            original_section_outline['num_chars'],
            raw_section_outline['location']['width'],
            raw_section_outline['location']['height']
        )
        if new_section_outline['suggestions'] != NOT_CHANGE:
            suggestion_flag = True
    if 'subsections' in original_section_outline:
        for k, v in original_section_outline['subsections'].items():
            if 'num_chars' in v:
                new_section_outline['subsections'][k]['suggestion'] = check_length(
                    result_json['subsections'][k]['description'], 
                    v['num_chars'],
                    raw_section_outline['subsections'][k]['location']['width'],
                    raw_section_outline['subsections'][k]['location']['height']
                )
                if new_section_outline['subsections'][k]['suggestion'] != NOT_CHANGE:
                    suggestion_flag = True

    return new_section_outline, suggestion_flag

def get_img_ratio(img_path):
    img = Image.open(img_path)
    return {
        'width': img.width,
        'height': img.height
    }

def get_img_ratio_in_section(content_json):
    res = {}
    if 'path' in content_json:
        res[content_json['path']] = get_img_ratio(content_json['path'])

    if 'subsections' in content_json:
        for subsection_name, val in content_json['subsections'].items():
            if 'path' in val:
                res[val['path']] = get_img_ratio(val['path'])

    return res


def get_snapshot_from_section(leaf_section, section_name, name_to_hierarchy, leaf_name, section_code, empty_poster_path='poster.pptx'):
    hierarchy = name_to_hierarchy[leaf_name]
    hierarchy_overflow_name = f'tmp/overflow_check_<{section_name}>_<{leaf_section}>_hierarchy_{hierarchy}'
    run_code_with_utils(section_code, utils_functions)
    poster = Presentation(empty_poster_path)
    # add border regardless of the hierarchy
    curr_location = add_border_hierarchy(
        poster, 
        name_to_hierarchy, 
        hierarchy, 
        border_width=10,
        # regardless=True
    )
    if not leaf_section in curr_location:
        leaf_section = section_name
    save_presentation(poster, file_name=f"{hierarchy_overflow_name}.pptx")
    ppt_to_images(
        f"{hierarchy_overflow_name}.pptx", 
        hierarchy_overflow_name, 
        dpi=200
    )
    poster_image_path = os.path.join(f"{hierarchy_overflow_name}", "slide_0001.jpg")
    poster_image = Image.open(poster_image_path)

    poster_width = emu_to_inches(poster.slide_width)
    poster_height = emu_to_inches(poster.slide_height)
    locations = convert_pptx_bboxes_json_to_image_json(
        curr_location, 
        poster_width, 
        poster_height
    )
    zoomed_in_img = zoom_in_image_by_bbox(
        poster_image, 
        locations[leaf_name], 
        padding=0.01
    )
    # save the zoomed_in_img
    zoomed_in_img.save(f"{hierarchy_overflow_name}_zoomed_in.jpg")
    return curr_location, zoomed_in_img, f"{hierarchy_overflow_name}_zoomed_in.jpg"